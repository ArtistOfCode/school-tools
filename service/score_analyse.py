import logging
import os
from typing import List

import numpy as np
from openpyxl import load_workbook
from openpyxl.workbook import Workbook
from openpyxl.worksheet.worksheet import Worksheet
from pptx import Presentation
from pptx.util import Inches

from model.score_model import ClassScore, SubjectScore
from model.student_model import is_valid_stu, Student, STU_DTYPE
from model.subject_model import Subjects
from service.excel_styles import set_cell, set_title_cell, set_float_cell, CellIndex, set_center_cell

DATA_FILES = ['一年级.xlsx', '二年级.xlsx', '三年级.xlsx', '四年级.xlsx', '五年级.xlsx', '六年级.xlsx']


class ScoreAnalyseService:

    def __init__(self, root_dir):
        self.file_paths = [f'{root_dir}/data/read/{f}' for f in DATA_FILES]
        self.result_path = f'{root_dir}/data/成绩分析结果.xlsx'
        self.ppt_template_path = f'{root_dir}/data/成绩分析模板.pptx'
        self.ppt_result_path = f'{root_dir}/data/成绩分析结果.pptx'

    # 校级分析
    def school_analyse(self):
        school_workbook: Workbook = Workbook()
        school_workbook.remove(school_workbook['Sheet'])
        school_ppt = Presentation(self.ppt_template_path)
        for file_path in self.file_paths:
            title = f'{os.path.basename(file_path)}'.replace('.xlsx', '')

            workbook: Workbook = load_workbook(file_path, True, False, True)
            school_score = self.grade_analyse(title, workbook)
            grade_sheet: Worksheet = school_workbook.create_sheet(title)

            grade_layout = school_ppt.slide_layouts[1]
            grade_slide = school_ppt.slides.add_slide(grade_layout)
            grade_slide.shapes.title.text = f'{title}成绩分析'

            row = CellIndex()
            for subject in Subjects: self.write_class(grade_sheet, school_score, subject, row)
            column = CellIndex(12)
            for subject in Subjects: self.write_care_stu(grade_sheet, school_score, subject, CellIndex(), column)
            for subject in Subjects: self.write_pptx(school_ppt, title, school_score, subject)

            workbook.close()
            logging.info(f'{title}分析完成！')
        school_workbook.save(self.result_path)
        school_ppt.save(self.ppt_result_path)
        logging.info('分析结果保存完成！')

    # 年级分析
    def grade_analyse(self, grade_name, workbook: Workbook):
        school_score: List[ClassScore] = []
        # 第一遍循环分析基本指标和一类关爱指标
        for class_score in self.class_analyse(grade_name, workbook):
            class_score.analyse()
            school_score.append(class_score)
        # 第二遍循环分析二类关爱指标
        for class_score in school_score:
            class_score.analyse_care(school_score[-1])
        return school_score

    # 班级分析
    @staticmethod
    def class_analyse(grade_name, workbook: Workbook):
        grade_students = []
        for sheetname in workbook.sheetnames:
            sheet = workbook[sheetname]
            class_students = [Student(row).to_tuple for row in sheet.iter_rows(min_row=2) if is_valid_stu(row)]
            grade_students.extend(class_students)
            yield ClassScore(grade_name, sheetname, np.array(class_students, dtype=STU_DTYPE))
        yield ClassScore(grade_name, '校平', np.array(grade_students, dtype=STU_DTYPE))

    # 保存分析结果
    @staticmethod
    def write_class(grade_sheet: Worksheet, school_score: List[ClassScore], subject: Subjects, row: CellIndex):
        is_low_grade = (grade_sheet.title == '一年级' or grade_sheet.title == '二年级')
        if is_low_grade and subject == Subjects.ENGLISH: return
        subject_code, subject_name = subject.value

        care_score = getattr(school_score[0], subject_code).care_stu_2[0]
        pass_name = '三科' if not is_low_grade and subject == Subjects.TWO else ''
        care_name = f'率({care_score})' if is_low_grade else '平均分'

        headers = ['班级', '总人数', '平均分', f'{pass_name}及格人数', f'{pass_name}及格率', '特优人数', '特优率',
                   f'关爱{care_name}']

        set_title_cell(grade_sheet.cell(row.value, 1), subject_name)
        row.next()

        for idx, header in enumerate(headers): set_title_cell(grade_sheet.cell(row.value, idx + 1), header)
        row.next()

        for idx, class_score in enumerate(school_score):
            _subject: SubjectScore = getattr(class_score, subject_code)
            row_index = row.value
            set_cell(grade_sheet.cell(row_index, 1), class_score.name)
            set_cell(grade_sheet.cell(row_index, 2), class_score.total_stu)
            set_float_cell(grade_sheet.cell(row_index, 3), _subject.mean)
            set_cell(grade_sheet.cell(row_index, 4), _subject.pass_stu[0])
            set_float_cell(grade_sheet.cell(row_index, 5), _subject.pass_stu[1])
            if subject != Subjects.ENGLISH:
                set_cell(grade_sheet.cell(row_index, 6), _subject.top_stu[0])
                set_float_cell(grade_sheet.cell(row_index, 7), _subject.top_stu[1])
            if is_low_grade:
                set_float_cell(grade_sheet.cell(row_index, 8), _subject.care_stu_2[2])
            else:
                set_float_cell(grade_sheet.cell(row_index, 8), _subject.care_stu_1[1])
            row.next()
        row.next()

    # 保存关爱学生
    @staticmethod
    def write_care_stu(grade_sheet: Worksheet, school_score: List[ClassScore], subject: Subjects, row: CellIndex,
                       column: CellIndex):
        is_low_grade = grade_sheet.title == '一年级' or grade_sheet.title == '二年级'
        if is_low_grade and subject == Subjects.ENGLISH: return
        subject_code, subject_name = subject.value

        name_col, score_col = column.value, column.value + 1
        column.next(3)

        for class_score in school_score:
            _subject: SubjectScore = getattr(class_score, subject_code)
            if _subject.care_stu_array.size == 0 or class_score.name == '校平': continue

            set_title_cell(grade_sheet.cell(row.value, name_col),
                           f'{class_score.name}班{subject_name}（{_subject.care_stu_array.size}）')
            row.next()

            if subject != Subjects.TWO:
                set_title_cell(grade_sheet.cell(row.value, name_col), '姓名')
                set_title_cell(grade_sheet.cell(row.value, score_col), '分数')
            else:
                set_title_cell(grade_sheet.cell(row.value, name_col), '姓名')
                set_title_cell(grade_sheet.cell(row.value, score_col), Subjects.CHINESE.value[1])
                set_title_cell(grade_sheet.cell(row.value, score_col + 1), Subjects.MATH.value[1])
                if class_score.is_low_grade:
                    set_title_cell(grade_sheet.cell(row.value, score_col + 2), Subjects.TWO.value[1])
                else:
                    set_title_cell(grade_sheet.cell(row.value, score_col + 2), Subjects.ENGLISH.value[1])
                    set_title_cell(grade_sheet.cell(row.value, score_col + 3), Subjects.TWO.value[1])
            row.next()

            for stu in _subject.care_stu_array:
                if subject != Subjects.TWO:
                    set_cell(grade_sheet.cell(row.value, name_col), stu['name'])
                    set_cell(grade_sheet.cell(row.value, score_col), stu[subject_code])
                else:
                    set_cell(grade_sheet.cell(row.value, name_col), stu['name'])
                    set_cell(grade_sheet.cell(row.value, score_col), stu[Subjects.CHINESE.value[0]])
                    set_cell(grade_sheet.cell(row.value, score_col + 1), stu[Subjects.MATH.value[0]])
                    if class_score.is_low_grade:
                        set_cell(grade_sheet.cell(row.value, score_col + 2), stu[Subjects.TWO.value[0]])
                    else:
                        set_cell(grade_sheet.cell(row.value, score_col + 2), stu[Subjects.ENGLISH.value[0]])
                        set_cell(grade_sheet.cell(row.value, score_col + 3), stu[Subjects.TWO.value[0]])
                row.next()
            row.next()

    # 保存分析结果PPT
    def write_pptx(self, school_ppt, title, school_score, subject: Subjects):
        is_low_grade = title == '一年级' or title == '二年级'
        if is_low_grade and subject == Subjects.ENGLISH: return
        subject_code, subject_name = subject.value

        # 添加幻灯片
        subject_layout = school_ppt.slide_layouts[2]
        class_layout = school_ppt.slide_layouts[3]
        subject_slide = school_ppt.slides.add_slide(subject_layout)
        subject_slide.shapes.title.text = f'{subject_name}情况分析'
        class_slide = school_ppt.slides.add_slide(class_layout)
        class_slide.shapes.title.text = f'{subject_name}情况分析'

        # 计算成绩表格表头
        if is_low_grade:
            care_score = getattr(school_score[0], subject_code).care_stu_2[0]
            if subject == Subjects.CHINESE or subject == Subjects.MATH:
                headers = ['班级', '平均分', '及格率', f'关爱率\v{care_score}', '总评', '与校\v平差', '与区\v平差',
                           '名次', '教者']
            else:
                headers = ['班级', '平均分', '及格人数', '及格率', f'关爱率\v{care_score}', '总评', '与校\v平差',
                           '与区\v平差', '名次', '班主任']
        else:
            if subject == Subjects.CHINESE or subject == Subjects.MATH:
                headers = ['班级', '平均分', '及格率', '关爱\v平均分', '特优率', '总评', '与校\v平差', '与区\v平差',
                           '名次', '教者']
            elif subject == Subjects.ENGLISH:
                headers = ['班级', '平均分', '及格率', f'关爱\v平均分', '总评', '与校\v平差', '与区\v平差', '名次',
                           '教者']
            else:
                headers = ['班级', '平均分', '三科\v及格人数', f'三科\v及格率', '关爱\v平均分', '总评', '与校\v平差',
                           '与区\v平差', '名次', '班主任']

        # 成绩表格排版
        ppt_width = school_ppt.slide_width.inches
        # ppt_height = school_ppt.slide_height.inches
        max_row = len(school_score) + 2
        max_column = len(headers)
        width = Inches(1.2)
        height = Inches(0.5)
        left = Inches((ppt_width - len(headers) * 1.2) / 2)
        top = Inches(1.5)
        table = class_slide.shapes.add_table(max_row, max_column, left, top, width, height).table

        for idx, header in enumerate(headers):
            table.columns[idx].width = width
            set_center_cell(table.cell(0, idx), header)

        row = CellIndex()

        color = None
        for idx, class_score in enumerate(school_score):
            _subject: SubjectScore = getattr(class_score, subject_code)
            row_idx = row.value
            table.rows[row_idx].height = height

            if is_low_grade:
                if subject == Subjects.CHINESE or subject == Subjects.MATH:
                    set_center_cell(table.cell(row_idx, 0), class_score.name, color)
                    set_center_cell(table.cell(row_idx, 1), self.to_string(_subject.mean), color)
                    set_center_cell(table.cell(row_idx, 2), self.to_string(_subject.pass_stu[1]), color)
                    set_center_cell(table.cell(row_idx, 3), self.to_string(_subject.care_stu_2[2]), color)
                else:
                    set_center_cell(table.cell(row_idx, 0), class_score.name, color)
                    set_center_cell(table.cell(row_idx, 1), self.to_string(_subject.mean), color)
                    set_center_cell(table.cell(row_idx, 2), self.to_string(_subject.pass_stu[0]), color)
                    set_center_cell(table.cell(row_idx, 3), self.to_string(_subject.pass_stu[1]), color)
                    set_center_cell(table.cell(row_idx, 4), self.to_string(_subject.care_stu_2[2]), color)
            else:
                if subject == Subjects.CHINESE or subject == Subjects.MATH:
                    set_center_cell(table.cell(row_idx, 0), class_score.name, color)
                    set_center_cell(table.cell(row_idx, 1), self.to_string(_subject.mean), color)
                    set_center_cell(table.cell(row_idx, 2), self.to_string(_subject.pass_stu[1]), color)
                    set_center_cell(table.cell(row_idx, 3), self.to_string(_subject.care_stu_1[1]), color)
                    set_center_cell(table.cell(row_idx, 4), self.to_string(_subject.top_stu[1]), color)
                elif subject == Subjects.ENGLISH:
                    set_center_cell(table.cell(row_idx, 0), class_score.name, color)
                    set_center_cell(table.cell(row_idx, 1), self.to_string(_subject.mean), color)
                    set_center_cell(table.cell(row_idx, 2), self.to_string(_subject.pass_stu[1]), color)
                    set_center_cell(table.cell(row_idx, 3), self.to_string(_subject.care_stu_1[1]), color)
                else:
                    set_center_cell(table.cell(row_idx, 0), class_score.name, color)
                    set_center_cell(table.cell(row_idx, 1), self.to_string(_subject.mean), color)
                    set_center_cell(table.cell(row_idx, 2), self.to_string(_subject.pass_stu[0]), color)
                    set_center_cell(table.cell(row_idx, 3), self.to_string(_subject.pass_stu[1]), color)
                    set_center_cell(table.cell(row_idx, 4), self.to_string(_subject.care_stu_1[1]), color)
            row.next()
        set_center_cell(table.cell(row.value, 0), '区平', color)
        row.next()

    @staticmethod
    def to_string(number):
        return str(round(number, 2))
